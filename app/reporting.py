from datetime import datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


INK = RGBColor(11, 31, 63)
MUTED = RGBColor(91, 111, 143)
BLUE = RGBColor(31, 96, 226)
BLUE_DARK = RGBColor(20, 69, 173)
BLUE_SOFT = RGBColor(232, 241, 255)
GREEN = RGBColor(8, 145, 96)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
LINE = RGBColor(214, 226, 242)
SURFACE = RGBColor(248, 251, 255)
WHITE = RGBColor(255, 255, 255)


def _text(value: Any, fallback: str = "n/a") -> str:
    if value is None or value == "":
        return fallback
    return str(value)


def _number(value: Any, fallback: str = "n/a") -> str:
    if isinstance(value, bool):
        return str(value)
    if isinstance(value, int):
        return f"{value:,}"
    if isinstance(value, float):
        if abs(value) < 1:
            return f"{value:.4f}"
        return f"{value:,.2f}".rstrip("0").rstrip(".")
    return _text(value, fallback)


def _short(value: Any, limit: int = 90) -> str:
    raw = _text(value, "")
    return raw if len(raw) <= limit else raw[: limit - 3].rstrip() + "..."


def _quality_color(label: Any) -> RGBColor:
    normalized = str(label or "").lower()
    if normalized in {"excellent", "good", "ready", "strong"}:
        return GREEN
    if normalized in {"fair", "needs_attention", "warning", "usable"}:
        return AMBER
    if normalized in {"poor", "weak", "blocked", "risk"}:
        return RED
    return BLUE


def _set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SURFACE


def _add_text(
    slide,
    value: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN | None = None,
):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = _text(value, "")
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_heading(slide, title: str, subtitle: str | None = None) -> None:
    _add_text(slide, "Automate Studio", 0.55, 0.22, 2.1, 0.28, 10, True, BLUE)
    _add_text(slide, title, 0.55, 0.55, 9.5, 0.45, 24, True, INK)
    if subtitle:
        _add_text(slide, subtitle, 0.57, 1.02, 10.4, 0.35, 10, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.38), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def _add_footer(slide, model_id: str | None, index: int) -> None:
    footer = f"Model ID: {_short(model_id, 26) if model_id else 'n/a'}"
    _add_text(slide, footer, 0.58, 7.05, 5.0, 0.2, 8, False, MUTED)
    _add_text(slide, f"{index}/6", 12.2, 7.05, 0.6, 0.2, 8, True, MUTED, PP_ALIGN.RIGHT)


def _add_card(
    slide,
    title: str,
    value: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    accent: RGBColor = BLUE,
    note: str | None = None,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    _add_text(slide, title.upper(), left + 0.2, top + 0.18, width - 0.35, 0.22, 8, True, MUTED)
    _add_text(slide, _short(value, 42), left + 0.2, top + 0.5, width - 0.35, 0.36, 17, True, INK)
    if note:
        _add_text(slide, _short(note, 58), left + 0.2, top + 0.9, width - 0.35, 0.24, 8, False, MUTED)


def _add_bullets(slide, lines: list[Any], left: float, top: float, width: float, height: float, size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    clean_lines = [line for line in lines if _text(line, "")]
    if not clean_lines:
        clean_lines = ["No saved details are available for this section yet."]

    for index, line in enumerate(clean_lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {_short(line, 150)}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(4)


def _add_section_box(slide, title: str, left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(1)
    _add_text(slide, title, left + 0.2, top + 0.18, width - 0.4, 0.32, 14, True, INK)


def _top_leaderboard_rows(record: dict, limit: int = 5) -> list[dict]:
    rows = record.get("leaderboard_snapshot") or []
    if rows:
        return rows[:limit]

    return [
        {
            "model_name": record.get("best_model"),
            "rank_score": record.get("rank_score"),
            "quality_label": record.get("quality_label"),
            "metrics": record.get("holdout_metrics") or {},
            "cross_validation": {"metric": record.get("primary_metric")},
            "status": "saved",
        }
    ]


def _top_features(record: dict, limit: int = 6) -> list[dict]:
    studio = record.get("explainability_studio") or {}
    features = studio.get("features") or record.get("top_features") or []
    return features[:limit]


def _metric_pairs(metrics: dict, limit: int = 6) -> list[tuple[str, Any]]:
    return [(key.replace("_", " ").title(), value) for key, value in list((metrics or {}).items())[:limit]]


def _history_lines(history: dict | None) -> list[str]:
    if not history:
        return []
    summary = history.get("summary") or {}
    best = summary.get("best_run") or {}
    latest = summary.get("latest_run") or {}
    lines = [
        f"Saved runs for this dataset/target: {_number(summary.get('filtered_runs'))}",
        f"Improved runs in view: {_number(summary.get('improved_runs'))}",
    ]
    if best:
        lines.append(f"Best saved run: {_text(best.get('best_model'))} at {_number(best.get('rank_score'))}")
    if latest and latest.get("score_delta") is not None:
        lines.append(f"Latest score delta: {_number(latest.get('score_delta'))}")
    return lines


def _build_title_slide(prs: Presentation, record: dict, dataset: dict | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)

    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BLUE
    ribbon.line.fill.background()

    _add_text(slide, "Automate Studio", 0.72, 0.72, 3.0, 0.35, 13, True, BLUE)
    _add_text(slide, "AutoML Project Report", 0.72, 1.2, 7.4, 0.7, 32, True, INK)
    _add_text(
        slide,
        f"{_text(record.get('best_model'), 'Best model')} for {_text(record.get('target_column'), 'selected target')}",
        0.75,
        2.0,
        8.0,
        0.45,
        15,
        False,
        MUTED,
    )

    generated = datetime.now().strftime("%d %b %Y, %I:%M %p")
    _add_card(slide, "Dataset", dataset.get("name") if dataset else record.get("dataset_id"), 0.75, 3.02, 2.85, 1.25, BLUE)
    _add_card(slide, "Problem", record.get("problem_type"), 3.9, 3.02, 2.35, 1.25, BLUE)
    _add_card(
        slide,
        _text(record.get("primary_metric"), "Metric"),
        _number(record.get("rank_score")),
        6.55,
        3.02,
        2.3,
        1.25,
        _quality_color(record.get("quality_label")),
        record.get("quality_label"),
    )
    _add_card(slide, "Generated", generated, 9.15, 3.02, 3.05, 1.25, BLUE)

    _add_section_box(slide, "Executive Summary", 0.75, 4.7, 11.45, 1.55)
    summary_lines = [
        f"Best model: {_text(record.get('best_model'))}",
        f"Target: {_text(record.get('target_column'))}",
        f"Candidates trained: {_number(record.get('trained_models'))} of {_number(record.get('candidate_models'))}",
        f"Prediction endpoint: {_text(record.get('prediction_api'))}",
    ]
    _add_bullets(slide, summary_lines, 1.0, 5.18, 10.9, 0.82, 11)
    _add_footer(slide, record.get("model_id"), 1)


def _build_readiness_slide(prs: Presentation, record: dict, dataset: dict | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_heading(slide, "Dataset And Target Readiness", "What the selected dataset/model run used for training.")

    _add_card(slide, "Rows", dataset.get("rows") if dataset else None, 0.65, 1.75, 1.75, 1.0, BLUE)
    _add_card(slide, "Columns", dataset.get("columns") if dataset else None, 2.6, 1.75, 1.75, 1.0, BLUE)
    _add_card(slide, "Raw features", record.get("raw_feature_count"), 4.55, 1.75, 1.75, 1.0, BLUE)
    _add_card(slide, "Model features", record.get("model_feature_count"), 6.5, 1.75, 1.75, 1.0, BLUE)
    _add_card(slide, "Dropped", record.get("dropped_feature_count"), 8.45, 1.75, 1.75, 1.0, AMBER)
    _add_card(slide, "Leakage flags", record.get("leakage_feature_count"), 10.4, 1.75, 1.75, 1.0, RED)

    _add_section_box(slide, "Training Context", 0.65, 3.12, 5.95, 2.8)
    context = record.get("run_context") or {}
    context_lines = [
        f"Dataset ID: {_text(record.get('dataset_id'))}",
        f"Target column: {_text(record.get('target_column'))}",
        f"Problem type: {_text(record.get('problem_type'))}",
        f"Run type: {_text(context.get('mode'), 'standard training')}",
        f"Selected features: {_number(record.get('selected_feature_count'))}",
        f"Excluded features: {_number(record.get('excluded_feature_count'))}",
    ]
    _add_bullets(slide, context_lines, 0.9, 3.62, 5.45, 1.85, 10)

    _add_section_box(slide, "Readiness Notes", 6.9, 3.12, 5.55, 2.8)
    leakage = record.get("leakage_features") or []
    class_balance = record.get("class_balance") or {}
    notes = []
    if leakage:
        notes.append(f"Auto-excluded possible leakage: {', '.join(_short(item.get('name'), 20) for item in leakage[:3])}")
    if class_balance.get("available"):
        notes.append(
            f"Class balance: {class_balance.get('severity')} with minority share {_number(class_balance.get('minority_share'))}"
        )
    notes.extend((record.get("next_actions") or [])[:3])
    _add_bullets(slide, notes, 7.15, 3.62, 5.05, 1.85, 10)
    _add_footer(slide, record.get("model_id"), 2)


def _build_leaderboard_slide(prs: Presentation, record: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_heading(slide, "Model Leaderboard", "Every report is tied to the selected trained model snapshot.")

    rows = _top_leaderboard_rows(record)
    table_shape = slide.shapes.add_table(len(rows) + 1, 5, Inches(0.65), Inches(1.75), Inches(12.0), Inches(3.1))
    table = table_shape.table
    table.columns[0].width = Inches(0.7)
    table.columns[1].width = Inches(4.6)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(2.9)

    headers = ["Rank", "Model", "Score", "Quality", "Validation"]
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)

    for row_index, item in enumerate(rows, start=1):
        cv = item.get("cross_validation") or {}
        validation = "No CV saved"
        if cv.get("available") or cv.get("mean") is not None:
            validation = f"{_text(cv.get('metric'), record.get('primary_metric'))}: {_number(cv.get('mean'))}"
            if cv.get("std") is not None:
                validation += f" +/- {_number(cv.get('std'))}"

        values = [
            row_index,
            _short(item.get("model_name"), 48),
            _number(item.get("rank_score")),
            _text(item.get("quality_label") or record.get("quality_label")),
            validation,
        ]
        for column, value in enumerate(values):
            cell = table.cell(row_index, column)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = INK

    _add_section_box(slide, "Run Coverage", 0.65, 5.25, 5.75, 1.25)
    coverage = [
        f"Candidate models: {_number(record.get('candidate_models'))}",
        f"Successfully trained: {_number(record.get('trained_models'))}",
        f"Failed models: {_number(record.get('failed_models'))}",
    ]
    _add_bullets(slide, coverage, 0.9, 5.72, 5.2, 0.55, 10)

    _add_section_box(slide, "Selected Winner", 6.7, 5.25, 5.95, 1.25)
    winner = [
        f"Best model: {_text(record.get('best_model'))}",
        f"Primary metric: {_text(record.get('primary_metric'))}",
        f"Saved API: {_text(record.get('prediction_api'))}",
    ]
    _add_bullets(slide, winner, 6.95, 5.72, 5.45, 0.55, 10)
    _add_footer(slide, record.get("model_id"), 3)


def _build_accuracy_slide(prs: Presentation, record: dict, history: dict | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_heading(slide, "Accuracy And Improvement Notes", "Baseline, holdout scores, tuning, and saved experiment context.")

    metrics = _metric_pairs(record.get("holdout_metrics") or {}, 6)
    if not metrics:
        metrics = [(record.get("primary_metric") or "Rank score", record.get("rank_score"))]
    for index, (label, value) in enumerate(metrics[:4]):
        _add_card(slide, label, _number(value), 0.65 + (index * 3.0), 1.72, 2.65, 1.0, BLUE)

    _add_section_box(slide, "Baseline Comparison", 0.65, 3.1, 3.9, 2.65)
    baseline = record.get("baseline_metrics") or {}
    baseline_lines = [f"{label}: {_number(value)}" for label, value in _metric_pairs(baseline, 5)]
    _add_bullets(slide, baseline_lines, 0.9, 3.58, 3.4, 1.58, 10)

    _add_section_box(slide, "Accuracy Booster", 4.82, 3.1, 3.9, 2.65)
    booster = record.get("accuracy_improver") or {}
    booster_lines = []
    if booster.get("summary"):
        booster_lines.append(booster.get("summary"))
    booster_lines.extend(item.get("detail") for item in (booster.get("checks") or [])[:3])
    booster_lines.extend(booster.get("actions") or [])
    _add_bullets(slide, booster_lines[:5], 5.07, 3.58, 3.4, 1.58, 9)

    _add_section_box(slide, "Experiment History", 8.98, 3.1, 3.65, 2.65)
    lines = _history_lines(history)
    tuning = record.get("tuning_studio") or {}
    if tuning:
        lines.append(f"Tuned models: {_number(tuning.get('tuned_count'))}")
        lines.append(f"Improved tuned models: {_number(tuning.get('improved_count'))}")
    _add_bullets(slide, lines[:5], 9.23, 3.58, 3.15, 1.58, 9)
    _add_footer(slide, record.get("model_id"), 4)


def _build_explainability_slide(prs: Presentation, record: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_heading(slide, "Explainability Summary", "Top feature drivers saved with this model run.")

    features = _top_features(record)
    _add_section_box(slide, "Top Drivers", 0.65, 1.7, 6.15, 4.75)
    if features:
        max_importance = max(float(item.get("importance") or 0) for item in features) or 1
        for index, item in enumerate(features):
            y = 2.24 + (index * 0.58)
            name = _short(item.get("feature"), 38)
            importance = float(item.get("importance") or 0)
            _add_text(slide, name, 0.95, y, 2.6, 0.22, 9, True, INK)
            track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.65), Inches(y + 0.04), Inches(2.35), Inches(0.12))
            track.fill.solid()
            track.fill.fore_color.rgb = BLUE_SOFT
            track.line.fill.background()
            bar_width = max(0.04, 2.35 * min(importance / max_importance, 1))
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.65), Inches(y + 0.04), Inches(bar_width), Inches(0.12))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BLUE
            bar.line.fill.background()
            _add_text(slide, _number(importance), 6.05, y, 0.55, 0.2, 8, False, MUTED, PP_ALIGN.RIGHT)
    else:
        _add_bullets(slide, ["Train a fresh model to save feature importance details."], 0.95, 2.3, 5.4, 1.0, 11)

    _add_section_box(slide, "Interpretation", 7.1, 1.7, 5.55, 4.75)
    studio = record.get("explainability_studio") or {}
    interpretation = []
    if studio.get("summary"):
        interpretation.append(studio.get("summary"))
    for feature in features[:4]:
        feature_name = _text(feature.get("feature"))
        detail = feature.get("interpretation") or feature.get("impact") or feature.get("strength")
        interpretation.append(f"{feature_name}: {_text(detail)}")
    _add_bullets(slide, interpretation, 7.35, 2.18, 5.05, 2.95, 10)
    _add_footer(slide, record.get("model_id"), 5)


def _build_deployment_slide(prs: Presentation, record: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_heading(slide, "Deployment And Next Actions", "How to use the selected model and what to improve next.")

    _add_card(slide, "Prediction API", record.get("prediction_api"), 0.65, 1.72, 5.65, 1.15, BLUE)
    _add_card(slide, "Threshold", _number(record.get("recommended_threshold")), 6.6, 1.72, 2.3, 1.15, BLUE)
    _add_card(slide, "Created", record.get("created_at"), 9.2, 1.72, 3.1, 1.15, BLUE)

    _add_section_box(slide, "Recommended Next Work", 0.65, 3.22, 5.75, 2.75)
    _add_bullets(slide, (record.get("next_actions") or [])[:6], 0.9, 3.72, 5.2, 1.75, 10)

    _add_section_box(slide, "Operational Checklist", 6.7, 3.22, 5.65, 2.75)
    checklist = [
        "Download and version the model artifact before sharing.",
        "Use batch scoring for CSV validation before production use.",
        "Review low-confidence predictions manually.",
        "Retrain after adding new rows or fixing weak-class recall.",
        "Export leaderboard and explainability CSVs for audit notes.",
    ]
    _add_bullets(slide, checklist, 6.95, 3.72, 5.15, 1.75, 10)
    _add_footer(slide, record.get("model_id"), 6)


def build_model_report_pptx(record: dict, dataset: dict | None = None, history: dict | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _build_title_slide(prs, record, dataset)
    _build_readiness_slide(prs, record, dataset)
    _build_leaderboard_slide(prs, record)
    _build_accuracy_slide(prs, record, history)
    _build_explainability_slide(prs, record)
    _build_deployment_slide(prs, record)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
